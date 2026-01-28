import json        # データを保存する「JSON形式」を扱うためのツール
import sys         # システム（PC）の情報を扱うためのツール
import time        # 待ち時間を設定するためのツール
import os          # ファイルやフォルダの操作をするためのツール
import re          # 特定の文字パターン（正規表現）を探すためのツール
import ctypes      # Windowsの高度な設定（タスクバーのアイコン表示など）に使うツール
import pickle      # Pythonのデータをそのまま保存・復元するためのツール
from datetime import datetime  # 日時（いま何時？）を扱うためのツール
from pathlib import Path       # フォルダの住所（パス）を賢く扱うためのツール

# --- 外部ライブラリ（後からインストールした特別な道具箱） ---
import google.generativeai as genai  # GoogleのAI「Gemini」を使うための道具
from PyQt6.QtWidgets import (       # アプリの「ボタン」や「入力欄」などの部品
    QApplication, QComboBox, QLabel, QListWidget, QMainWindow, QMenu,
    QMessageBox, QPushButton, QStyle, QSystemTrayIcon, QVBoxLayout,
    QHBoxLayout, QWidget, QLineEdit, QListWidgetItem, QProgressDialog, QInputDialog
)
from PyQt6.QtGui import QAction, QColor, QIcon  # アイコンや色、メニューの動作を司る部品
from PyQt6.QtCore import Qt, QSize              # 位置やサイズなどの細かな設定用の部品
from docx import Document           # Wordファイルを作るための道具
from openpyxl import Workbook       # Excelファイルを作るための道具
from pptx import Presentation       # PowerPointファイルを作るための道具
from PIL import Image               # 画像ファイルを扱うための道具

# --- Google Drive API 関連（ネット上のドライブへ送るための設定） ---
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request

# --- パス（住所）・環境設定 ---
# 実行している場所が「exeファイル」の中か、普通の「プログラム」の中かで住所を切り替えています
if getattr(sys, 'frozen', False):
    BASE_DIR = Path(sys.executable).parent
else:
    BASE_DIR = Path(__file__).resolve().parent

# データの保存先（dataフォルダなど）を指定しています
DATA_DIR = BASE_DIR / "data"
HISTORY_PATH = DATA_DIR / "history.jsonl"
CONFIG_PATH = DATA_DIR / "config.json"
ICON_PATH = BASE_DIR / "icon.ico"
CRED_PATH = BASE_DIR / "credentials.json"  # Google Cloudから取得した「合鍵」の場所
TOKEN_PATH = DATA_DIR / "token.pickle"      # ログイン情報を「一時保存」する場所
WINDOWS_FORBIDDEN_CHARS = re.compile(r'[\\/:*?"<>|]') # ファイル名に使えない禁止文字のリスト

# Google Driveで「自分のアプリが作ったファイルだけ」を操作する権限を設定
SCOPES = ['https://www.googleapis.com/auth/drive.file']

# --- Google Drive 認証・アップロードエンジン ---

def get_drive_service(parent):
    """Google Driveを使うための「通行許可」をもらう処理"""
    creds = None
    # 以前のログイン情報（トークン）があればそれを使う
    if TOKEN_PATH.exists():
        with open(TOKEN_PATH, 'rb') as token:
            creds = pickle.load(token)
    
    # ログイン情報がない、または期限切れなら新しくログインする
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        else:
            if not CRED_PATH.exists():
                QMessageBox.critical(parent, "認証エラー", "credentials.json がありません。")
                return None
            flow = InstalledAppFlow.from_client_secrets_file(str(CRED_PATH), SCOPES)
            creds = flow.run_local_server(port=0)
        
        # ログイン情報を保存して、次回から入力を省けるようにする
        DATA_DIR.mkdir(parents=True, exist_ok=True)
        with open(TOKEN_PATH, 'wb') as token:
            pickle.dump(creds, token)
            
    return build('drive', 'v3', credentials=creds)

def upload_to_drive(file_path, mime_type, parent):
    """作成したファイルをGoogleドライブへ「送信」する処理"""
    service = get_drive_service(parent)
    if not service: return None
    
    file_metadata = {'name': file_path.name} # 送るファイルの名前
    media = MediaFileUpload(str(file_path), mimetype=mime_type) # ファイルの本体
    
    try:
        # ドライブ上にファイルを作成
        file = service.files().create(body=file_metadata, media_body=media, fields='id').execute()
        return file.get('id')
    except Exception as e:
        print(f"Drive Upload Error: {e}")
        return None

# --- アプリの見た目（デザイン）の設定：QSS ---
MODERN_STYLE = """
QMainWindow { background-color: #1e1e1e; } /* 全体の背景は暗いグレー */
QWidget { color: #e0e0e0; font-family: "Segoe UI", "Meiryo"; font-size: 14px; } /* 文字は明るいグレー */
/* ボタンや入力欄の角を丸くしたり、色をつけたりしています */
QPushButton#primary { background-color: #007acc; color: white; }
QPushButton#danger { color: #ff6b6b; }
"""

# --- 便利な小道具関数（ユーティリティ） ---

def append_history_entry(entry):
    """「何を実行したか」の履歴をファイルに1行ずつ書き足す処理"""
    DATA_DIR.mkdir(parents=True, exist_ok=True)
    with HISTORY_PATH.open("a", encoding="utf-8", newline="\n") as h:
        h.write(json.dumps(entry, ensure_ascii=False) + "\n")

def load_run_records(limit=10):
    """過去の実行データを読み込んでリストにする処理（最新10件）"""
    runs_dir = DATA_DIR / "runs"
    if not runs_dir.exists(): return []
    records = []
    for d in runs_dir.iterdir():
        p = d / "run.json"
        if p.exists():
            try:
                with p.open("r", encoding="utf-8") as h: records.append(json.load(h))
            except: continue
    records.sort(key=lambda x: x.get("created_at", ""), reverse=True)
    return records[:limit]

def save_run_record(record):
    """今回の実行内容を「個別フォルダ」の中に保存する処理"""
    run_id = record.get("run_id")
    if not run_id: return False
    run_dir = DATA_DIR / "runs" / run_id
    run_dir.mkdir(parents=True, exist_ok=True)
    try:
        with (run_dir / "run.json").open("w", encoding="utf-8") as h:
            json.dump(record, h, ensure_ascii=False, indent=2)
        return True
    except: return False

def sanitize_output_basename(raw):
    """ファイル名として使えない文字を「_」に置き換える処理"""
    return WINDOWS_FORBIDDEN_CHARS.sub("_", (raw or "").strip()) or "result"

def get_api_key(parent=None):
    """Geminiを使うための「APIキー」を読み込む。なければ入力画面を出す処理"""
    if CONFIG_PATH.exists():
        with open(CONFIG_PATH, "r", encoding="utf-8") as f:
            c = json.load(f)
            if "api_key" in c: return c["api_key"]
    # キーを隠しながら入力してもらうためのダイアログを表示
    k, ok = QInputDialog.getText(parent, "Gemini APIキー", "キーを入力:", QLineEdit.EchoMode.Password)
    if ok and k:
        DATA_DIR.mkdir(parents=True, exist_ok=True)
        with open(CONFIG_PATH, "w", encoding="utf-8") as f: json.dump({"api_key": k}, f)
        return k
    return None

# --- AI解析のメイン処理（ここがアプリの心臓部） ---
def execute_gemini_process(record, owner):
    """AIに指示を出し、ファイルを作り、ドライブへ送る一連の大きな流れ"""
    api_key = get_api_key(owner)
    if not api_key: return
    
    # 処理中にユーザーが操作できないように「待機中」のバーを出す
    p = QProgressDialog("AI解析中 ＆ ドライブ同期中...", "キャンセル", 0, 0, owner)
    p.setWindowModality(Qt.WindowModality.WindowModal)
    p.show(); QApplication.processEvents()
    
    try:
        # AIの設定と呼び出し
        genai.configure(api_key=api_key)
        model = genai.GenerativeModel('models/gemini-2.5-flash')
        run_id, doc_format, purpose = record["run_id"], record.get("doc_format"), record["purpose"]
        
        # AIへの「お願い（プロンプト）」を組み立てる
        prompt = f"目的: {purpose}\n"
        if doc_format == "Word": prompt += "指示: 構造化レポート形式で出力。"
        elif doc_format == "Excel": prompt += "指示: タブ区切りの表形式で出力。"
        elif doc_format == "PowerPoint": prompt += "指示: スライド構成案を出力。"
        else: prompt += "指示: 装飾記号を使わないプレーンテキストで出力。"
        
        # 保存してあるスクリーンショット画像も一緒にAIに送る準備
        prompt_parts = [prompt]
        for img in record.get("captures", []):
            img_p = DATA_DIR / "captures" / run_id / img
            if img_p.exists(): prompt_parts.append(Image.open(img_p))
        
        # AIに送信して結果（テキスト）を受け取る
        res = model.generate_content(prompt_parts)
        ai_text = res.text
        
        base = record["output_basename"]
        out_d = DATA_DIR / "outputs" / run_id; out_d.mkdir(parents=True, exist_ok=True)
        
        m_type = "text/plain" # ネットに送る際の「ファイルの種類」のメモ

        # 選んだ形式に合わせて、WordやExcelのファイルとして書き出す
        if doc_format == "Word":
            target = out_d / f"{base}.docx"
            m_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            doc = Document(); doc.add_heading(purpose, 0); doc.add_paragraph(ai_text); doc.save(str(target))
        elif doc_format == "Excel":
            target = out_d / f"{base}.xlsx"
            m_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            wb = Workbook(); wb.active["A1"] = ai_text; wb.save(str(target))
        elif doc_format == "PowerPoint":
            target = out_d / f"{base}.pptx"
            m_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            prs = Presentation(); s = prs.slides.add_slide(prs.slide_layouts[1]); s.shapes.title.text = purpose; s.placeholders[1].text = ai_text; prs.save(str(target))
        else:
            target = out_d / f"{base}.txt"
            with target.open("w", encoding="utf-8") as h: h.write(ai_text.replace("**", "").replace("#", "").strip())
        
        # 実行完了の記録を付け、ファイルを開く
        record["status"] = "done"; record["output_file"] = target.name
        save_run_record(record); os.startfile(str(target))

        # --- Googleドライブへ自動アップロード ---
        drive_id = upload_to_drive(target, m_type, owner)
        
        if drive_id:
            QMessageBox.information(owner, "成功", f"解析完了！ドライブにも保存しました。")

    except Exception as e:
        QMessageBox.critical(owner, "エラー", str(e))
    finally:
        p.close() # 待機バーを閉じる

# --- ウィンドウクラス（アプリの各「画面」の設計図） ---

class NewRunWindow(QMainWindow):
    """「新しくAIに頼む」ときの入力画面"""
    def __init__(self, run_id):
        super().__init__(); self.run_id = run_id; self.capture_index = 1
        self.created_at = datetime.now().isoformat(timespec="seconds")
        self.setWindowTitle("新規実行 - SnipAI")
        if ICON_PATH.exists(): self.setWindowIcon(QIcon(str(ICON_PATH)))
        self._build_ui()
        
    def _build_ui(self):
        """画面上にボタンやリストなどの部品を並べる処理"""
        c = QWidget(); l = QVBoxLayout(c); l.setContentsMargins(20, 20, 20, 20); l.setSpacing(12)
        self.output_name_input = QLineEdit(); self.output_name_input.setPlaceholderText("ファイル名を入力")
        self.purpose_combo = QComboBox(); self.purpose_combo.addItems(["まとめる", "解説する", "ドキュメント化"])
        self.doc_format_combo = QComboBox(); self.doc_format_combo.addItems(["Word", "Excel", "PowerPoint"])
        self.doc_format_combo.setEnabled(False) # 最初はWordとかは選べないようにしておく
        # 「ドキュメント化」を選んだときだけWordなどの形式を選べるように切り替える
        self.purpose_combo.currentTextChanged.connect(lambda t: self.doc_format_combo.setEnabled(t == "ドキュメント化"))
        self.capture_list = QListWidget()
        bl = QHBoxLayout(); cb = QPushButton("📸 スクショ追加"); db = QPushButton("🗑 選択削除")
        db.setObjectName("danger"); bl.addWidget(cb); bl.addWidget(db)
        rb = QPushButton("🚀 AI解析を実行"); rb.setObjectName("primary"); sb = QPushButton("💾 保存")
        
        # 上から順番に部品を詰め込んでいく（レイアウト）
        l.addWidget(QLabel("出力ファイル名")); l.addWidget(self.output_name_input)
        row = QHBoxLayout(); v1, v2 = QVBoxLayout(), QVBoxLayout()
        v1.addWidget(QLabel("目的")); v1.addWidget(self.purpose_combo)
        v2.addWidget(QLabel("形式")); v2.addWidget(self.doc_format_combo)
        row.addLayout(v1); row.addLayout(v2); l.addLayout(row)
        l.addWidget(QLabel("スクリーンショット")); l.addWidget(self.capture_list)
        l.addLayout(bl); l.addSpacing(15); l.addWidget(rb); l.addWidget(sb)
        
        # ボタンを押した時にどの処理を呼ぶか（コネクト）
        cb.clicked.connect(self.capture_screenshot); db.clicked.connect(self.remove_selected)
        rb.clicked.connect(self.run_now); sb.clicked.connect(self.save_only)
        self.setCentralWidget(c); self.resize(450, 600)
        
    def _collect(self):
        """入力された情報をまとめて「実行データ」という1つのセットにする処理"""
        base = sanitize_output_basename(self.output_name_input.text())
        fmt = self.doc_format_combo.currentText() if self.doc_format_combo.isEnabled() else "Text"
        ext = {"Word": ".docx", "Excel": ".xlsx", "PowerPoint": ".pptx", "Text": ".txt"}[fmt]
        return {
            "run_id": self.run_id, "created_at": self.created_at, "purpose": self.purpose_combo.currentText(),
            "doc_format": fmt if fmt != "Text" else None, "save_target": "local", 
            "captures": [self.capture_list.item(i).text() for i in range(self.capture_list.count())],
            "status": "ready", "output_basename": base, "output_file": f"{base}{ext}"
        }
        
    def capture_screenshot(self):
        """画面を一時的に隠して、今のウィンドウをパシャリと撮る処理"""
        self.hide(); time.sleep(0.3)
        try:
            import mss, mss.tools, win32gui
            d = DATA_DIR / "captures" / self.run_id; d.mkdir(parents=True, exist_ok=True)
            f = f"{self.capture_index:03d}.png"; fp = d / f
            # いま一番手前にあるウィンドウの範囲を特定して撮影
            h = win32gui.GetForegroundWindow(); l, t, r, b = win32gui.GetWindowRect(h)
            with mss.mss() as sct:
                img = sct.grab({"left": l, "top": t, "width": r-l, "height": b-t})
                mss.tools.to_png(img.rgb, img.size, output=str(fp))
            self.capture_list.addItem(f); self.capture_index += 1
        finally: self.show()
        
    def remove_selected(self):
        """リストで選んでいるスクショを削除する処理"""
        for i in self.capture_list.selectedItems(): self.capture_list.takeItem(self.capture_list.row(i))
        
    def run_now(self):
        """「今すぐ実行」ボタンが押されたときの動き"""
        if self.capture_list.count() == 0: return QMessageBox.warning(self, "SnipAI", "スクショを撮ってください")
        d = self._collect(); save_run_record(d); append_history_entry({**d, "id": self.run_id})
        execute_gemini_process(d, self); self.close()
        
    def save_only(self):
        """「保存だけ」して後で実行する場合の動き"""
        d = self._collect(); save_run_record(d); append_history_entry({**d, "id": self.run_id}); self.close()

class HistoryWindow(QMainWindow):
    """「過去の履歴」を一覧表示する画面"""
    def __init__(self):
        super().__init__(); self.setWindowTitle("実行履歴")
        self.list_widget = QListWidget(); self.list_widget.setSpacing(6)
        if ICON_PATH.exists(): self.setWindowIcon(QIcon(str(ICON_PATH)))
        c = QWidget(); l = QVBoxLayout(c); l.addWidget(self.list_widget); self.setCentralWidget(c); self.resize(500, 450)
        
    def refresh(self):
        """履歴ファイルを読み直して、画面を最新にする処理"""
        self.list_widget.clear()
        for r in load_run_records():
            item = QListWidgetItem(); w = QWidget(); lay = QHBoxLayout(w)
            lay.setContentsMargins(10, 2, 10, 2)
            lbl = QLabel(f"{r['created_at'][5:16].replace('T', ' ')} | {r['purpose']}")
            btn = QPushButton("📄 開く" if r.get('status') == "done" else "▶ 実行")
            btn.setFixedWidth(110); btn.setFixedHeight(34)
            btn.clicked.connect(lambda _, rec=r: self._handle(rec))
            lay.addWidget(lbl, 1); lay.addWidget(btn); item.setSizeHint(QSize(100, 58))
            self.list_widget.addItem(item); self.list_widget.setItemWidget(item, w)
            
    def _handle(self, r):
        """履歴のボタンが押されたとき「ファイルを開く」か「今からAIに投げる」か判断する処理"""
        if r['status'] == "done":
            p = DATA_DIR / "outputs" / r["run_id"] / r["output_file"]
            if p.exists(): os.startfile(str(p))
            else: QMessageBox.warning(self, "エラー", "ファイルなし")
        else: execute_gemini_process(r, self); self.refresh()

class HomeWindow(QMainWindow):
    """アプリを起動して最初に表示されるメインメニュー画面"""
    def __init__(self):
        super().__init__(); self.setWindowTitle("SnipAI"); self.setStyleSheet(MODERN_STYLE)
        if ICON_PATH.exists(): self.setWindowIcon(QIcon(str(ICON_PATH)))
        c = QWidget(); l = QVBoxLayout(c); l.setContentsMargins(30, 30, 30, 30); l.setSpacing(15)
        b1 = QPushButton("🆕 新規実行"); b1.setObjectName("primary"); b1.setFixedHeight(50)
        b2 = QPushButton("📜 履歴を表示"); b2.setFixedHeight(50)
        b1.clicked.connect(self.open_new); b2.clicked.connect(self.open_hist)
        l.addWidget(b1); l.addWidget(b2); self.setCentralWidget(c); self.resize(350, 250)
        
    def open_new(self):
        """新規実行画面を開く"""
        rid = datetime.now().strftime("%Y%m%d_%H%M%S"); self.nw = NewRunWindow(rid); self.nw.setStyleSheet(MODERN_STYLE); self.nw.show()
        
    def open_hist(self):
        """履歴画面を開く"""
        if not hasattr(self, 'hw'): self.hw = HistoryWindow(); self.hw.setStyleSheet(MODERN_STYLE)
        self.hw.refresh(); self.hw.show()
        
    def closeEvent(self, e): 
        """「×」ボタンを押しても完全に終了せず、タスクトレイに隠れるようにする設定"""
        self.hide(); e.ignore()

def main():
    """アプリの起動そのものを行う一番外側の処理"""
    try:
        # タスクバーでアイコンを正しく表示させるためのWindows用の呪文
        myappid = 'mickyyya.snipai.v1' 
        ctypes.windll.shell32.SetCurrentProcessExplicitAppUserModelID(myappid)
    except: pass

    app = QApplication(sys.argv); app.setQuitOnLastWindowClosed(False)
    icon = QIcon(str(ICON_PATH)) if ICON_PATH.exists() else app.style().standardIcon(QStyle.StandardPixmap.SP_ComputerIcon)
    home = HomeWindow(); home.setWindowIcon(icon)
    
    # タスクトレイ（画面右下の小さなアイコン集）に常駐させる設定
    tray = QSystemTrayIcon(icon, app); menu = QMenu()
    menu.addAction("ホーム").triggered.connect(home.show); menu.addAction("終了").triggered.connect(app.quit)
    tray.setContextMenu(menu); tray.show(); home.show()
    
    return app.exec() # アプリが動いている間、ずっとここにとどまる

# プログラムが直接実行されたときだけ、main()を呼び出す
if __name__ == "__main__": main()